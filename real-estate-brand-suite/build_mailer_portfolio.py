"""Luxury 9x6" landscape direct-mail postcard template — Prestige Properties.

Two slides, 9" wide x 6" tall:
  Slide 1 — Front: full-bleed hero photo + address overlay band + Offered At
  Slide 2 — Back : photo strip + narrative + specs + agent + CTA

Brand-matched to the 11x17 bi-fold brochure (same palette and typography)
so the print suite ships as one cohesive system.
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_DEEP = RGBColor(0xA8, 0x89, 0x47)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CHARCOAL = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF5, 0xEF, 0xE2)
CREAM_DEEP = RGBColor(0xEA, 0xE0, 0xCB)
GRAY = RGBColor(0x88, 0x84, 0x7A)

SERIF = "Playfair Display"
SANS = "Montserrat"


def rect(slide, x, y, w, h, fill=None, line=None, line_pt=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None: s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_pt)
    return s


def hline(slide, x1, y, x2, color=GOLD, pt=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = color; c.line.width = Pt(pt)
    return c


def text(slide, x, y, w, h, content, *, font=SERIF, size=12, bold=False, italic=False,
         color=BLACK, align="left", anchor="top", line_spacing=1.15, tracking=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    lines = content.split("\n") if isinstance(content, str) else content
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
        if tracking is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def photo_placeholder(slide, x, y, w, h, label="PHOTO", sub=""):
    rect(slide, x, y, w, h, fill=CREAM, line=CREAM_DEEP, line_pt=0.75)
    inset = 0.06
    rect(slide, x + inset, y + inset, w - 2 * inset, h - 2 * inset, fill=None, line=GOLD, line_pt=0.5)
    cx, cy = x + w / 2, y + h / 2 - 0.10
    d = 0.2
    dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    dia.shadow.inherit = False
    dia.fill.solid(); dia.fill.fore_color.rgb = GOLD
    dia.line.fill.background()
    text(slide, x, y + h / 2 + 0.02, w, 0.28, label, font=SANS, size=9, bold=True, color=GOLD_DEEP, align="center", tracking=500)
    if sub:
        text(slide, x, y + h / 2 + 0.32, w, 0.22, sub, font=SANS, size=7, italic=True, color=GRAY, align="center")


def build_front(slide):
    W, H = 9.0, 6.0
    # full-bleed hero across left 2/3
    photo_placeholder(slide, 0, 0, 6.0, H, label="HERO LISTING PHOTO", sub="Full-bleed print-ready image")
    # right 1/3 cream panel
    rect(slide, 6.0, 0, 3.0, H, fill=WHITE)
    # gold vertical accent where panels meet
    rect(slide, 6.0, 0, 0.04, H, fill=GOLD)
    # top gold strip across right panel
    rect(slide, 6.0, 0, 3.0, 0.22, fill=GOLD)
    # wordmark
    text(slide, 6.0, 0.45, 3.0, 0.3, "PRESTIGE PROPERTIES",
         font=SANS, size=10, bold=True, color=BLACK, align="center", tracking=600)
    text(slide, 6.0, 0.80, 3.0, 0.22, "Personal Service. Prestige Results.",
         font=SERIF, size=9, italic=True, color=GRAY, align="center")
    # hairline
    hline(slide, 6.5, 1.15, 8.5, color=GOLD, pt=0.6)
    # JUST LISTED eyebrow
    text(slide, 6.0, 1.35, 3.0, 0.3, "JUST LISTED", font=SANS, size=11, bold=True,
         color=GOLD_DEEP, align="center", tracking=800)
    # hairline
    hline(slide, 6.8, 1.75, 8.2, color=GOLD, pt=1.0)
    # address
    text(slide, 6.0, 1.95, 3.0, 0.75, "[STREET\nADDRESS]", font=SERIF, size=19,
         color=BLACK, align="center", line_spacing=1.0)
    text(slide, 6.0, 2.9, 3.0, 0.35, "[CITY, STATE ZIP]", font=SERIF, size=11,
         italic=True, color=CHARCOAL, align="center")
    # price block
    rect(slide, 6.2, 3.5, 2.6, 0.5, fill=BLACK)
    text(slide, 6.2, 3.56, 2.6, 0.4, "OFFERED AT  $[PRICE]",
         font=SANS, size=9, bold=True, color=GOLD, align="center", anchor="middle", tracking=400)
    # specs mini-strip
    text(slide, 6.0, 4.2, 3.0, 0.3, "[#] BED  ·  [#] BATH  ·  [SQ FT]",
         font=SANS, size=9, color=CHARCOAL, align="center", tracking=300)
    hline(slide, 6.8, 4.6, 8.2, color=CREAM_DEEP, pt=0.6)
    # CTA
    rect(slide, 6.3, 4.8, 2.4, 0.55, fill=None, line=GOLD, line_pt=1.0)
    text(slide, 6.3, 4.88, 2.4, 0.4, "VIEW THE LISTING",
         font=SANS, size=9, bold=True, color=BLACK, align="center", anchor="middle", tracking=500)
    # bottom gold strip
    rect(slide, 6.0, 5.78, 3.0, 0.22, fill=GOLD)


def build_back(slide):
    W, H = 9.0, 6.0
    # top gold banner
    rect(slide, 0, 0, W, 0.22, fill=GOLD)
    # 3-photo strip
    gap = 0.12
    m = 0.35
    gw = (W - 2 * m - 2 * gap) / 3
    gh = 2.2
    photo_placeholder(slide, m, 0.5, gw, gh, label="PHOTO 01")
    photo_placeholder(slide, m + gw + gap, 0.5, gw, gh, label="PHOTO 02")
    photo_placeholder(slide, m + 2 * (gw + gap), 0.5, gw, gh, label="PHOTO 03")
    # narrative section
    hline(slide, m, 2.95, m + 1.2, color=GOLD, pt=1.0)
    text(slide, m, 3.05, W - 2 * m, 0.3, "THE RESIDENCE", font=SANS, size=10,
         bold=True, color=GOLD_DEEP, tracking=700)
    text(slide, m, 3.40, W - 2 * m, 0.45, "A Sanctuary of Timeless Elegance",
         font=SERIF, size=18, italic=True, color=BLACK)
    text(slide, m, 3.95, W - 2 * m, 0.9,
         "Set on a private drive in one of the Southeast Region's most coveted "
         "Premier Neighborhoods, this custom-built estate blends architectural heritage "
         "with effortless modern living.",
         font=SANS, size=9, color=CHARCOAL, line_spacing=1.45)
    # specs black strip
    rect(slide, m, 4.85, W - 2 * m, 0.38, fill=BLACK)
    text(slide, m, 4.88, W - 2 * m, 0.32,
         "[#] BEDROOMS   ·   [#] BATHS   ·   [SQ FT]   ·   [ACRES]   ·   BUILT [YEAR]",
         font=SANS, size=9, bold=True, color=GOLD, align="center", anchor="middle", tracking=350)
    # agent block bottom
    # left: agent block
    text(slide, m, 5.30, 4.0, 0.3, "EXCLUSIVELY REPRESENTED BY",
         font=SANS, size=8, color=GOLD_DEEP, tracking=600)
    text(slide, m, 5.55, 4.0, 0.3, "Your Name Here",
         font=SERIF, size=14, color=BLACK)
    # right: contact + CTA
    text(slide, W - m - 4.0, 5.30, 4.0, 0.3, "+1 (555) 000-0000   ·   info@prestigeproperties.com",
         font=SANS, size=9, color=CHARCOAL, align="right")
    text(slide, W - m - 4.0, 5.55, 4.0, 0.3, "PRESTIGE PROPERTIES",
         font=SANS, size=10, bold=True, color=BLACK, align="right", tracking=500)


def build(out: Path):
    prs = Presentation()
    prs.slide_width = Inches(9.0)
    prs.slide_height = Inches(6.0)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    rect(s1, 0, 0, 9, 6, fill=WHITE)
    build_front(s1)
    s1.notes_slide.notes_text_frame.text = (
        "MAILER FRONT. 9x6 inch landscape direct-mail postcard. Print on 14pt cover "
        "or 130 lb silk with 0.125 inch bleed on each side (final trim 9 x 6). "
        "Replace hero photo by right-click > Change Picture on the cream rectangle."
    )

    s2 = prs.slides.add_slide(blank)
    rect(s2, 0, 0, 9, 6, fill=WHITE)
    build_back(s2)
    s2.notes_slide.notes_text_frame.text = (
        "MAILER BACK. USPS side A typically reserved for address block + indicia. "
        "If mailing via EDDM or first class, reserve the right third for address "
        "panel or swap this design to back-as-address-side. Reach-out details go "
        "in agent block."
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="./Prestige_Properties_Mailer_Portfolio.pptx")
    args = ap.parse_args()
    p = build(Path(args.out))
    print(f"WROTE {p}")
