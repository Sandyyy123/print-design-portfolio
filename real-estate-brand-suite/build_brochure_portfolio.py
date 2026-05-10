"""Luxury 11x17 bi-fold real estate brochure template — Prestige Properties.

Outputs an editable .pptx with two 17"x11" slides:
  Slide 1 — outside spread: [back cover | front cover]
  Slide 2 — inside spread : [inside left | inside right]

Panels are 8.5"x11" each. Fold guide drawn as faint center line.
Photo placeholders = cream rectangles with centered label; client can
right-click > Change Picture, or delete + Insert > Picture.
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

GOLD = RGBColor(0xC9, 0xA9, 0x61)
GOLD_DEEP = RGBColor(0xA8, 0x89, 0x47)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
CHARCOAL = RGBColor(0x2B, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF5, 0xEF, 0xE2)
CREAM_DEEP = RGBColor(0xEA, 0xE0, 0xCB)
GRAY = RGBColor(0x88, 0x84, 0x7A)

SERIF = "Playfair Display"
SANS = "Montserrat"


def rect(slide, x, y, w, h, fill=None, line=None, line_pt=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_pt)
    return s


def hline(slide, x1, y, x2, color=GOLD, pt=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    return c


def vline(slide, x, y1, y2, color=GOLD, pt=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(pt)
    return c


def text(
    slide,
    x,
    y,
    w,
    h,
    content,
    *,
    font=SERIF,
    size=12,
    bold=False,
    italic=False,
    color=BLACK,
    align="left",
    anchor="top",
    line_spacing=1.15,
    tracking=None,  # hundredths of percent (e.g. 200 = +2% spacing)
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if tracking is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb


def photo_placeholder(slide, x, y, w, h, label="HERO PHOTO", sub="Right-click > Change Picture"):
    rect(slide, x, y, w, h, fill=CREAM, line=CREAM_DEEP, line_pt=0.75)
    # inner hairline border for elegance
    inset = 0.08
    rect(slide, x + inset, y + inset, w - 2 * inset, h - 2 * inset, fill=None, line=GOLD, line_pt=0.5)
    # camera glyph (simple diamond)
    cx, cy = x + w / 2, y + h / 2 - 0.15
    d = 0.25
    dia = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    dia.shadow.inherit = False
    dia.fill.solid()
    dia.fill.fore_color.rgb = GOLD
    dia.line.fill.background()
    text(slide, x, y + h / 2 + 0.05, w, 0.35, label, font=SANS, size=10, bold=True, color=GOLD_DEEP, align="center", tracking=400)
    text(slide, x, y + h / 2 + 0.40, w, 0.3, sub, font=SANS, size=7.5, italic=True, color=GRAY, align="center")


def build_back_cover(slide, x0):
    # left panel: x0=0, width=8.5
    W = 8.5
    # top gold thin banner
    rect(slide, x0, 0, W, 0.35, fill=GOLD)
    # agent photo
    photo_placeholder(slide, x0 + 2.5, 1.1, 3.5, 4.0, label="AGENT PHOTO", sub="Right-click > Change Picture")
    # agent name
    text(slide, x0 + 0.5, 5.3, W - 1, 0.6, "Your Name Here", font=SERIF, size=28, color=BLACK, align="center")
    # title row
    text(
        slide,
        x0 + 0.5,
        5.95,
        W - 1,
        0.3,
        "REALTOR®   ·   LUXURY SPECIALIST",
        font=SANS,
        size=9,
        color=GOLD_DEEP,
        align="center",
        tracking=500,
    )
    hline(slide, x0 + 3.25, 6.45, x0 + 5.25)
    # contact block
    text(
        slide,
        x0 + 0.5,
        6.7,
        W - 1,
        1.6,
        "+1 (555) 000-0000\ninfo@prestigeproperties.com\nwww.prestigeproperties.com",
        font=SANS,
        size=11,
        color=CHARCOAL,
        align="center",
        line_spacing=1.55,
    )
    hline(slide, x0 + 3.25, 8.55, x0 + 5.25)
    # tagline
    text(
        slide,
        x0 + 0.5,
        8.75,
        W - 1,
        0.6,
        "“Personal Service. Prestige Results.”",
        font=SERIF,
        size=15,
        italic=True,
        color=CHARCOAL,
        align="center",
    )
    # wordmark bottom
    text(slide, x0, 10.05, W, 0.3, "PRESTIGE PROPERTIES", font=SANS, size=11, bold=True, color=BLACK, align="center", tracking=600)
    text(slide, x0, 10.4, W, 0.25, "PREMIER NEIGHBORHOODS  ·  SOUTHEAST REGION", font=SANS, size=7.5, color=GRAY, align="center", tracking=400)
    # bottom gold banner
    rect(slide, x0, 10.8, W, 0.2, fill=GOLD)


def build_front_cover(slide, x0):
    W = 8.5
    # hero photo, full bleed on three sides, leaves bottom for title treatment
    photo_placeholder(slide, x0 + 0.4, 0.4, W - 0.8, 7.3, label="HERO PHOTO", sub="Full-bleed listing image")
    # presenting label
    text(
        slide,
        x0 + 0.5,
        7.95,
        W - 1,
        0.3,
        "PRESENTING",
        font=SANS,
        size=10,
        bold=True,
        color=GOLD_DEEP,
        align="center",
        tracking=800,
    )
    hline(slide, x0 + 3.75, 8.35, x0 + 4.75, color=GOLD, pt=1.0)
    # address lines
    text(
        slide,
        x0 + 0.5,
        8.5,
        W - 1,
        0.9,
        "1234 Elegance Lane",
        font=SERIF,
        size=32,
        color=BLACK,
        align="center",
    )
    text(
        slide,
        x0 + 0.5,
        9.35,
        W - 1,
        0.55,
        "Premier Neighborhoods, Southeast Region XXXXX",
        font=SERIF,
        size=18,
        italic=True,
        color=CHARCOAL,
        align="center",
    )
    hline(slide, x0 + 0.75, 9.95, x0 + W - 0.75, color=GOLD, pt=0.5)
    # offered at + wordmark
    text(slide, x0 + 0.75, 10.1, 3.5, 0.4, "OFFERED AT  $2,495,000", font=SANS, size=10, bold=True, color=CHARCOAL, align="left", tracking=400)
    text(slide, x0 + W - 4.25, 10.1, 3.5, 0.4, "PRESTIGE PROPERTIES", font=SANS, size=10, bold=True, color=GOLD_DEEP, align="right", tracking=600)


def build_inside_left(slide, x0):
    W = 8.5
    # large photo top
    photo_placeholder(slide, x0 + 0.5, 0.5, W - 1, 5.8, label="INTERIOR / EXTERIOR PHOTO", sub="Feature shot")
    # gold hairline
    hline(slide, x0 + 0.5, 6.65, x0 + 2.0, color=GOLD, pt=1.0)
    # eyebrow
    text(slide, x0 + 0.5, 6.75, W - 1, 0.3, "THE RESIDENCE", font=SANS, size=10, bold=True, color=GOLD_DEEP, tracking=700)
    # headline serif italic
    text(
        slide,
        x0 + 0.5,
        7.15,
        W - 1,
        1.0,
        "A Sanctuary of\nTimeless Elegance",
        font=SERIF,
        size=30,
        italic=True,
        color=BLACK,
        line_spacing=1.0,
    )
    # narrative
    narrative = (
        "Set on a private drive in one of the Southeast Region’s most coveted "
        "Premier Neighborhoods, this custom-built estate blends architectural heritage with "
        "effortless modern living. Sunlit interiors, hand-selected finishes, "
        "and seamless indoor–outdoor flow define a home designed for "
        "quiet luxury and spirited gatherings alike."
    )
    text(slide, x0 + 0.5, 8.55, W - 1, 1.9, narrative, font=SANS, size=11, color=CHARCOAL, line_spacing=1.55)
    # bottom specs strip on gold
    rect(slide, x0 + 0.5, 10.3, W - 1, 0.45, fill=BLACK)
    text(
        slide,
        x0 + 0.5,
        10.33,
        W - 1,
        0.4,
        "5 BEDROOMS   ·   6 BATHS   ·   6,420 SQ FT   ·   1.2 ACRES   ·   BUILT 2022",
        font=SANS,
        size=9.5,
        bold=True,
        color=GOLD,
        align="center",
        anchor="middle",
        tracking=400,
    )


def build_inside_right(slide, x0):
    W = 8.5
    # 2x2 photo grid
    margin = 0.5
    gap = 0.18
    gw = (W - 2 * margin - gap) / 2
    gh = 2.4
    y0 = 0.5
    photo_placeholder(slide, x0 + margin, y0, gw, gh, label="PHOTO 01", sub="")
    photo_placeholder(slide, x0 + margin + gw + gap, y0, gw, gh, label="PHOTO 02", sub="")
    photo_placeholder(slide, x0 + margin, y0 + gh + gap, gw, gh, label="PHOTO 03", sub="")
    photo_placeholder(slide, x0 + margin + gw + gap, y0 + gh + gap, gw, gh, label="PHOTO 04", sub="")
    # hairline + eyebrow
    hline(slide, x0 + 0.5, 5.9, x0 + 2.0, color=GOLD, pt=1.0)
    text(slide, x0 + 0.5, 6.0, W - 1, 0.3, "PRESTIGE FEATURES", font=SANS, size=10, bold=True, color=GOLD_DEEP, tracking=700)
    # features two columns
    col1 = (
        "•  Chef’s kitchen with La Cornue range\n"
        "•  Primary suite with spa bath & dual closets\n"
        "•  Wine cellar & tasting lounge\n"
        "•  Home theater with tiered seating"
    )
    col2 = (
        "•  Saltwater pool & covered loggia\n"
        "•  Four-car garage with EV charging\n"
        "•  Whole-home automation & security\n"
        "•  Detached guest casita"
    )
    text(slide, x0 + 0.5, 6.4, (W - 1) / 2, 2.2, col1, font=SANS, size=11, color=CHARCOAL, line_spacing=1.5)
    text(slide, x0 + 0.5 + (W - 1) / 2, 6.4, (W - 1) / 2, 2.2, col2, font=SANS, size=11, color=CHARCOAL, line_spacing=1.5)
    # divider
    hline(slide, x0 + 0.5, 8.85, x0 + W - 0.5, color=CREAM_DEEP, pt=0.5)
    # representation block
    text(slide, x0 + 0.5, 9.0, W - 1, 0.3, "EXCLUSIVELY REPRESENTED BY", font=SANS, size=9, color=GOLD_DEEP, align="center", tracking=700)
    text(slide, x0 + 0.5, 9.35, W - 1, 0.55, "Your Name Here", font=SERIF, size=22, color=BLACK, align="center")
    text(
        slide,
        x0 + 0.5,
        9.95,
        W - 1,
        0.3,
        "+1 (555) 000-0000   ·   info@prestigeproperties.com   ·   www.prestigeproperties.com",
        font=SANS,
        size=9.5,
        color=CHARCOAL,
        align="center",
    )
    hline(slide, x0 + 3.0, 10.35, x0 + 5.5, color=GOLD)
    text(
        slide,
        x0 + 0.5,
        10.45,
        W - 1,
        0.3,
        "PRESTIGE PROPERTIES",
        font=SANS,
        size=10,
        bold=True,
        color=BLACK,
        align="center",
        tracking=600,
    )


def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(17)
    prs.slide_height = Inches(11)
    blank = prs.slide_layouts[6]

    # ---- Slide 1: outside spread ----
    s1 = prs.slides.add_slide(blank)
    rect(s1, 0, 0, 17, 11, fill=WHITE)
    build_back_cover(s1, x0=0)
    build_front_cover(s1, x0=8.5)
    # fold guide (subtle dashed would be ideal; use very light gray thin)
    vline(s1, 8.5, 0, 11, color=CREAM_DEEP, pt=0.5)
    # trim/fold label (outside)
    text(s1, 8.25, 10.82, 0.5, 0.2, "│", font=SANS, size=6, color=GRAY, align="center")
    # notes
    notes = s1.notes_slide.notes_text_frame
    notes.text = (
        "OUTSIDE SPREAD — prints on one side of 11x17. Left panel = back cover. "
        "Right panel = front cover. Fold along the center vertical. "
        "To replace any photo: right-click the placeholder rectangle > Format Shape > "
        "Fill > Picture or texture fill > Insert > From File. Or delete the rectangle "
        "and Insert > Picture. Editable text fields: agent name, address, city, "
        "price, contact details."
    )

    # ---- Slide 2: inside spread ----
    s2 = prs.slides.add_slide(blank)
    rect(s2, 0, 0, 17, 11, fill=WHITE)
    build_inside_left(s2, x0=0)
    build_inside_right(s2, x0=8.5)
    vline(s2, 8.5, 0, 11, color=CREAM_DEEP, pt=0.5)
    notes2 = s2.notes_slide.notes_text_frame
    notes2.text = (
        "INSIDE SPREAD — prints on the reverse side of 11x17. Left panel = property "
        "narrative + hero interior. Right panel = 4-photo grid + features + agent "
        "representation block. Fold along the center vertical."
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="./Prestige_Properties_Brochure_Portfolio.pptx")
    args = ap.parse_args()
    path = build(Path(args.out))
    print(f"WROTE {path}")
